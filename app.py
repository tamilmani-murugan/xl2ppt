from pptx import Presentation
from pptx.chart.data import CategoryChartData
from openpyxl import load_workbook
from ast import literal_eval


def get_data_from_xl( sheet , xl_cord):
    for i, elem in enumerate(xl_cord):
        if isinstance(elem, str):
            xl_cord[i] = sheet[xl_cord[i]].value
        else:
            xl_cord[i] = get_data_from_xl(sheet , elem)
    return xl_cord

def get_categories_from_ppt(categories):
    if any(2 == len(x) for x in categories):
        md = {}
        for x in categories:
            if x[0] in md:
                md[x[0]].append(x[1])
            else:
                md[x[0]] = [x[1]]
        return ([[key] + [val] for key, val in md.items()])
    else:
        return ([x[0] for x in categories])

def get_data_from_ppt(prs , slide , shape):
    chart_data = []
    for plot in shape.chart.plots:
        plot_data = []
        plot_data.append(get_categories_from_ppt(plot.categories.flattened_labels))
        for series in  plot.series:
            series_data = []
            series_data.append(["{0:.0%}".format(s) if isinstance(s, float) else s for s in series.values])
            series_data.append(series.name)
            plot_data.append(series_data)
        chart_data.append(plot_data)
    return chart_data

def add_notes(slide,note) :
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = text_frame.text + '\n' + note

def validate_data(slide , chart_n , ppt_data , xl_data , column_n) :
    if len(ppt_data) != len(xl_data):
        note = f"Number of fields are not matching with excel for chart {chart_n}"
        add_notes(slide,note)
        return 1
    for i in range(0,column_n) :
        if ppt_data[i][1:] != xl_data[i][:column_n]:
            note = f"Values are not matching with excel for chart {chart_n}"
            add_notes(slide,note)
            return 1
    return 0

def append_sub_category( chart_data , categories ) :
    for category in categories:
        if isinstance(category, str):
            chart_data.add_category(category)
            continue
        else:
            myc = chart_data.add_category(category[0])
            for sub_category in category[1]:
                myc.add_sub_category(sub_category)

def append_series( chart_data , serieses ) :
    for series in serieses:
        chart_data.add_series(series[1],series[0],'0%')
        #chart_data.add_series(series[1],[None],'0%')

def update_chart_data(prs , shape , output_file , sheet , line) :
    xl_data = get_data_from_xl( sheet , literal_eval(line.split('|')[3]))
    chart_data = CategoryChartData()
    for xl_plot in xl_data:
        append_sub_category( chart_data , xl_plot[0] )
        append_series( chart_data , xl_plot[1] )
    shape.chart.replace_data(chart_data)
    prs.save(output_file)

def show_shapes(slide):
    for shape in slide.shapes:
        print(f'ID : {shape.shape_id} & name : {shape.name} & type : {shape.shape_type}')
 
def update_txt_box(prs , shape, output_file , sheet , line ):
    shape.text_frame.paragraphs[0].runs[0].text = str(sheet[line.split('|')[3]].value)
    prs.save(output_file)

def update_table(prs , shape, output_file , sheet , line ):
    xl_cord = literal_eval(line.split('|')[3])
    for v,row in enumerate(shape.table.rows):
        for h, cell in enumerate(row.cells):
            if sheet[xl_cord[v][h]].value != None:
                cell.text_frame.paragraphs[0].runs[0].text = str(sheet[xl_cord[v][h]].value)
    prs.save(output_file)

def get_table_data(shape):
    r = []
    for row in shape.table.rows:
        c = []
        for cell in row.cells:
            c.append(cell.text)
        r.append(c)
    return r

def xl_to_ppt (sample_layout , sample_data , inf_data , output_file) :
    prs = Presentation(sample_layout)
    with open(inf_data) as f:
        for line in f:
            if line[0] != '#':
                slide = prs.slides[int(line.split('|')[0])-1]
                shape = slide.shapes[int(line.split('|')[1])-1]
                sheet = load_workbook(sample_data).worksheets[int(line.split('|')[2])-1]
                if 'TEXT_BOX' in str(shape.shape_type):
                    update_txt_box(prs , shape, output_file , sheet , line )
                elif 'TABLE' in str(shape.shape_type):
                    update_table(prs , shape , output_file , sheet , line )
                elif 'CHART' in str(shape.shape_type):
                    update_chart_data(prs , shape , output_file , sheet , line )

def xtract_ppt(sample_layout,output_file):
    prs = Presentation(sample_layout)
    with open(output_file, "w") as file:
        for sl_n,slide in enumerate(prs.slides,start =1):
            #show_shapes(slide)
            for sh_n,shape in enumerate(slide.shapes,start =1):
                if 'TEXT_BOX' in str(shape.shape_type):
                    file.write(f'{sl_n}|{sh_n}||{shape.text}\n')
                elif 'TABLE' in str(shape.shape_type):
                    file.write(f'{sl_n}|{sh_n}||{get_table_data(shape)}\n')
                elif 'CHART' in str(shape.shape_type):
                    file.write(f'{sl_n}|{sh_n}||{get_data_from_ppt(prs , slide , shape)}\n')
                else:
                    file.write(f'#{sl_n}|{sh_n}||Un Classified\n')


if __name__ == '__main__' :
    sample_layout = "C:\\Users\\822385\\Documents\\Scripts\\xl2ppt\\Dummy Slides.pptx"
    sample_data = "C:\\Users\\822385\\Documents\\Scripts\\xl2ppt\\Dummy Output.xlsx"
    inf_data = "C:\\Users\\822385\\Documents\\Scripts\\xl2ppt\\xl2ppt.inf"
    output_file = "C:\\Users\\822385\\Documents\\Scripts\\xl2ppt\\output.pptx"
    #xl_to_ppt (sample_layout , sample_data , inf_data , output_file)
    for line in xtract_ppt(sample_layout):
        pass
        #print(line)
